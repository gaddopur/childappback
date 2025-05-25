"""
Document Summarization Tool

This script provides functionality for summarizing various document formats using a Generative AI model. 
Supports both command-line usage and programmatic integration with web services.

Supported Formats:
- PDF (Research papers, articles)
- DOCX/DOC (Word documents)
- PPTX/PPT (PowerPoint presentations)
- XLSX/XLS (Excel spreadsheets)
- TXT (Plain text files)
- MD (Markdown documents)
- HTML (Web pages)
- EPUB (E-books)

Features:
- Multi-format support with intelligent text extraction
- Secure document validation (path safety, type, size)
- Format-aware summarization prompts
- API key rotation with failover capabilities
- Asynchronous processing for web scalability

Usage Patterns:

Command Line:
--------------
python pdf_summarizer.py "/path/to/file" --retries 5

Programmatic API:
-----------------
from pdf_summarizer import PDFSummarizer

summarizer = PDFSummarizer()

# Synchronous usage
summary = summarizer.summarize("document.pdf")

# Asynchronous usage
async def process_file():
    summary = await summarizer.async_summarize("document.pdf")
    # Handle the summary...

Web Service Integration:
------------------------
from fastapi import FastAPI, UploadFile
from pdf_summarizer import PDFSummarizer
import tempfile

app = FastAPI()
summarizer = PDFSummarizer()

@app.post("/summarize")
async def summarize_pdf(file: UploadFile):
    with tempfile.NamedTemporaryFile(delete=True) as tmp:
        tmp.write(await file.read())
        return await summarizer.async_summarize(tmp.name)
"""

import os
import sys
import asyncio
import logging
import time
from pathlib import Path
import pymupdf
from typing import Optional, Tuple
from google import genai
from google.genai import types
from google.api_core import exceptions
from models.model_config_Document_summarizer import (
    MODEL_CONFIG,
    PROMPT_TEMPLATES,
    SUPPORTED_TYPES
)

sys.path.append(os.path.abspath(os.path.join(os.path.dirname(__file__), '..')))
from childappback import settings
from models.api_key_manager import APIKeyManager

# Configure logging to output to both a file and the console
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[
        logging.StreamHandler()  # Console output
    ]
)

class PDFValidationError(Exception):
    """Custom exception for PDF validation failures."""
    pass

class PDFSummarizer:
    """
    PDFSummarizer Class

    This class provides the core functionality to summarize PDF documents using a 
    Generative AI model. It includes features for:
    - Secure validation of input PDFs, including path sanitization, file type, and size checks.
    - Robust text extraction from PDF files.
    - Fail-safe integration with the Generative AI model, supporting API key management, retries, and logging.
    - Asynchronous support for use in web applications or scalable systems.

    Usage:
    Create an instance of PDFSummarizer and use the `summarize` method for synchronous 
    processing or `async_summarize` for asynchronous processing. The class is designed 
    for reuse in various contexts, including command-line tools, APIs, or backend services.
    """
    
    def __init__(self, api_key_manager: APIKeyManager = None):
        """
        Initialize the summarizer with API key management and configuration.
        :param api_key_manager: APIKeyManager instance for managing API keys.
        """
        self.api_key_manager = api_key_manager or APIKeyManager()
        # Load configuration from imported module
        self.model_name = MODEL_CONFIG["model_name"]
        self.generation_config = MODEL_CONFIG["generation_config"]
        self.max_file_size = MODEL_CONFIG["max_file_size"]

        # Configure path restrictions
        self.allowed_path_prefix = Path(
            settings.MEDIA_ROOT
        ) / Path(MODEL_CONFIG["allowed_path_prefix"]).resolve()        
        logging.info(f"Allowed path prefix set to: {self.allowed_path_prefix}")

    def _get_client(self, api_key: str):
        """
        Retrieve or create a generative model for a given API key.
        :param api_key: The API key to configure the model.
        :return: Configured generative model or None if failed.
        """
        try:            
            client = genai.Client(api_key=api_key)

            logging.info(f"Created new model instance for key ending ...{api_key[-4:]}")
            return client

        except Exception as e:
            logging.error(f"Client creation failed: {str(e)}")
            return None

    async def async_summarize(self, pdf_path: str) -> Optional[str]:
        """
        Asynchronous interface to summarize a PDF file.
        :param pdf_path: Path to the PDF file.
        :return: Generated summary or None if failed.
        """
        loop = asyncio.get_event_loop()
        try:
            # Run the synchronous summarize method in an executor
            return await loop.run_in_executor(None, self.summarize, pdf_path)
        except Exception as e:
            logging.error(f"Async processing failed: {str(e)}")
            return None

    def summarize(self, pdf_path: str, max_retries: int = 3) -> Optional[str]:
        """
        Generate a summary for a PDF file.
        :param pdf_path: Path to the PDF file.
        :param max_retries: Maximum number of retry attempts.
        :return: Generated summary or None if failed.
        """
        total_start = time.time()
        try:
            # Validation timing
            valid_start = time.time()
            valid, reason = self._validate_pdf(pdf_path)
            valid_duration = time.time() - valid_start
            if not valid:
                logging.error(f"Validation failed in {valid_duration:.2f}s: {reason}")
                return None

            # Processing
            process_start = time.time()
            result = self._process_document(pdf_path, max_retries)
            process_duration = time.time() - process_start

            logging.info(
                f"Total processing time: {time.time()-total_start:.2f}s | "
                f"Validation: {valid_duration:.2f}s | "
                f"Core processing: {process_duration:.2f}s"
            )
            return result
        
        except Exception as e:
            logging.error(f"Total failure after {time.time()-total_start:.2f}s: {str(e)}")
            return None

    def _process_document(self, pdf_path: str, max_retries: int) -> Optional[str]:
        """
        Process the PDF file to generate a summary with retries.
        :param pdf_path: Path to the PDF file.
        :param max_retries: Maximum number of retry attempts.
        :return: Generated summary or None if failed.
        """
        for attempt in range(1, max_retries + 1):
            attempt_start = time.time()
            api_key = self.api_key_manager.get_available_key()
            if not api_key:
                logging.warning(f"Attempt {attempt} failed: No keys available")
                time.sleep(2 ** attempt)  # Exponential backoff for retries
                continue

            try:
                # Extract text from the PDF
                extract_start = time.time()
                text, ext = self._extract_text(pdf_path)
                extract_duration = time.time() - extract_start
                if not text:
                    return None

                # Get the model for the API key
                model_start = time.time()
                client = self._get_client(api_key)
                model_duration = time.time() - model_start
                if not client:
                    continue

                # Generate the summary
                gen_start = time.time()
                generation_model_config = types.GenerateContentConfig(
                temperature=self.generation_config["temperature"],
                top_p=self.generation_config["top_p"],
                top_k=self.generation_config["top_k"],
                max_output_tokens=self.generation_config["max_output_tokens"],
                )

                if not text.strip():
                    logging.error("No text content extracted from document")
                    return None
                
                # Truncate text to 30k characters to avoid token limits
                processing_text = text[:30000].strip()
                if not processing_text:
                    return None

                prompt = f"{PROMPT_TEMPLATES.get(ext, PROMPT_TEMPLATES['default'])}\n{processing_text}"
                logging.debug(f"Sending prompt to LLM:\n{prompt}")
                response = ""
                for chunk in client.models.generate_content_stream(
                    model=self.model_name,
                    contents=[types.Content(role="user", parts=[types.Part.from_text(text=prompt)])],
                    config=generation_model_config,
                ):
                    if chunk.text:
                        response += chunk.text

                self.api_key_manager.update_key_status(api_key, success=True)
                logging.info("LLM response received successfully")
                
                gen_duration = time.time() - gen_start
                # Update key usage status
                self.api_key_manager.update_key_status(api_key, success=True)
                logging.info(
                    f"Attempt {attempt} success | "
                    f"Extract: {extract_duration:.2f}s | "
                    f"Model: {model_duration:.2f}s | "
                    f"Gen: {gen_duration:.2f}s"
                )
                return response

            except exceptions.ResourceExhausted:
                error_duration = time.time() - attempt_start
                logging.warning(f"Attempt {attempt} failed in {error_duration:.2f}s: {str(e)}")
                self._handle_error(api_key, "Rate limit exceeded", attempt, max_retries)
            except exceptions.GoogleAPIError as e:
                error_duration = time.time() - attempt_start
                logging.warning(f"Attempt {attempt} failed in {error_duration:.2f}s: {str(e)}")
                self._handle_error(api_key, f"API Error: {str(e)}", attempt, max_retries)
            except Exception as e:
                error_duration = time.time() - attempt_start
                logging.warning(f"Attempt {attempt} failed in {error_duration:.2f}s: {str(e)}")
                self._handle_error(api_key, f"Unexpected error: {str(e)}", attempt, max_retries)

        logging.error(f"Failed after {max_retries} attempts")
        return None

    def _validate_pdf(self, pdf_path: str) -> Tuple[bool, str]:
        """
        Validate the PDF file for path safety, existence, type, and size.
        :param pdf_path: Path to the PDF file.
        :return: Tuple with validation status and reason.
        """
        try:
            path = Path(settings.MEDIA_ROOT) / pdf_path
            ext = path.suffix[1:].lower()

            # Check supported types
            if ext not in SUPPORTED_TYPES:
                return False, f"Unsupported file type: {ext}"

            # Debugging logs for resolved paths
            logging.debug(f"Resolved input path: {path}")
            logging.debug(f"Allowed path prefix: {self.allowed_path_prefix}")

            # Ensure the file path is within the allowed directory
            if not path.is_relative_to(self.allowed_path_prefix):
                return False, (
                    f"Security restriction: PDF must be in {self.allowed_path_prefix}\n"
                    f"Attempted path: {path}"
                )
            
            # Check file existence and type
            if not path.exists():
                return False, "File not found"
            
            # Check file size
            file_size = path.stat().st_size
            if file_size > self.max_file_size:
                return False, f"File too large ({file_size/1024/1024:.1f}MB > {self.max_file_size/1024/1024}MB limit)"
                
            return True, "Validation passed"
            
        except Exception as e:
            return False, f"Validation error: {str(e)}"

    def _extract_text(self, file_path: str) -> Optional[str]:
        """
        Extract text from supported document types
        """
        path = Path(settings.MEDIA_ROOT) / file_path
        ext = path.suffix[1:].lower()
        
        try:
            if ext == 'pdf':
                return self._extract_pdf_text(path), ext
            elif ext == 'docx':
                return self._extract_docx_text(path), ext
            elif ext == 'pptx':
                return self._extract_pptx_text(path), ext
            elif ext == 'xlsx':
                return self._extract_xlsx_text(path), ext
            elif ext in ['txt', 'md']:
                return self._extract_plain_text(path), ext
            elif ext == 'html':
                return self._extract_html_text(path), ext
            elif ext == 'epub':
                return self._extract_epub_text(path), ext
            else:
                logging.error(f"Unsupported file type: {ext}")
                return None, ext
        except Exception as e:
            logging.error(f"Text extraction failed for {ext}: {str(e)}")
            return None

    def _extract_pdf_text(self, path: Path) -> str:
        """Extract text from PDF using pymupdf"""
        text = []
        with pymupdf.open(path) as doc:
            for page in doc.pages():
                text.append(page.get_text())
        return "\n".join(text)

    def _extract_docx_text(self, path: Path) -> str:
        """Extract text from Word documents"""
        from docx import Document
        doc = Document(path)
        return "\n".join([para.text for para in doc.paragraphs])

    def _extract_pptx_text(self, path: Path) -> str:
        """Extract text from PowerPoint presentations"""
        from pptx import Presentation
        prs = Presentation(path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)

    def _extract_xlsx_text(self, path: Path) -> str:
        """Extract meaningful data from Excel files"""
        from openpyxl import load_workbook
        try:
            wb = load_workbook(filename=path, read_only=True, data_only=True)
            text = []
            
            for sheet in wb:
                # Skip empty sheets
                if sheet.max_row == 0 or sheet.max_column == 0:
                    continue
                    
                text.append(f"\n\nSheet: {sheet.title}")
                
                # Get header row if exists
                headers = []
                if sheet.max_row > 1:
                    headers = [cell.value for cell in sheet[1] if cell.value]
                
                # Process rows with data
                for row in sheet.iter_rows(min_row=2 if headers else 1, values_only=True):
                    row_data = []
                    for idx, cell in enumerate(row):
                        if cell is None:
                            continue
                        # Add header label if available
                        if idx < len(headers):
                            row_data.append(f"{headers[idx]}: {cell}")
                        else:
                            row_data.append(str(cell))
                    if row_data:
                        text.append(" | ".join(row_data))
            
            return "\n".join(text) if text else "No readable data found in spreadsheet"
        
        except Exception as e:
            logging.error(f"Excel processing error: {str(e)}")
            return ""

    def _extract_plain_text(self, path: Path) -> str:
        """Extract text from plain text files"""
        with open(path, 'r', encoding='utf-8') as f:
            return f.read()

    def _extract_html_text(self, path: Path) -> str:
        """Extract main content from HTML files"""
        from bs4 import BeautifulSoup
        with open(path, 'r', encoding='utf-8') as f:
            soup = BeautifulSoup(f, 'html.parser')
            return soup.get_text()

    def _extract_epub_text(self, path: Path) -> str:
        """Extract text from EPUB eBooks"""
        from ebooklib import epub
        book = epub.read_epub(path)
        text = []
        for item in book.get_items():
            if item.get_type() == epub.ITEM_DOCUMENT:
                text.append(item.get_content().decode('utf-8'))
        return "\n".join(text)

    def _handle_error(self, api_key: str, message: str, attempt: int, max_attempts: int):
        """
        Handle errors during API calls or retries.
        :param api_key: The API key used in the failed attempt.
        :param message: Error message.
        :param attempt: Current attempt number.
        :param max_attempts: Total number of allowed attempts.
        """
        logging.warning(f"Attempt {attempt}/{max_attempts} failed: {message}")
        self.api_key_manager.update_key_status(api_key, success=False)
        time.sleep(1)  # Brief pause between retries

if __name__ == "__main__":
    import argparse

    def main():
        """
        Main function to execute the PDF summarization tool.
        """
        parser = argparse.ArgumentParser(description='PDF Summarization Tool')
        parser.add_argument('pdf_path', type=str, help='Path to PDF file')
        parser.add_argument('--retries', type=int, default=3, help='Number of retry attempts')
        args = parser.parse_args()

        summarizer = PDFSummarizer()

        summary = summarizer.summarize(args.pdf_path, args.retries)
        
        if summary:
            print("# Summary\n")
            print(summary)
            exit(0)
        else:
            logging.error("Failed to generate summary")
            exit(1)

    main()

